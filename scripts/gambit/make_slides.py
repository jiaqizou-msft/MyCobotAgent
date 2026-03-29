"""Generate a 2-slide deck for the Surface Laptop Robot project."""
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
import os

OUTPUT = os.path.join(os.path.dirname(__file__), "..", "..", "temp", "SurfaceLaptopRobot_Deck.pptx")
IMG_DIR = os.path.join(os.path.dirname(__file__), "..", "..", "temp")

# Image paths from user's attachments
IMG_CAMERAS = os.path.join(IMG_DIR, "camera_check", "composite_now.jpg")
IMG_ANNOTATE = os.path.join(os.path.dirname(__file__), "..", "..", "temp", "camera_check", "all_now.jpg")

prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)

# ── Colors ──
BG_DARK = RGBColor(0x1a, 0x1a, 0x2e)
BG_MED = RGBColor(0x16, 0x21, 0x3e)
ACCENT = RGBColor(0x00, 0xbc, 0xd4)
WHITE = RGBColor(0xff, 0xff, 0xff)
LIGHT_GRAY = RGBColor(0xcc, 0xcc, 0xcc)
ORANGE = RGBColor(0xff, 0xa5, 0x00)
GREEN = RGBColor(0x00, 0xe6, 0x76)


def add_bg(slide, color):
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = color


def add_text_box(slide, left, top, width, height, text, font_size=18,
                 color=WHITE, bold=False, alignment=PP_ALIGN.LEFT):
    txBox = slide.shapes.add_textbox(Inches(left), Inches(top), Inches(width), Inches(height))
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(font_size)
    p.font.color.rgb = color
    p.font.bold = bold
    p.alignment = alignment
    return tf


def add_bullet(tf, text, font_size=16, color=LIGHT_GRAY, bold=False, level=0):
    p = tf.add_paragraph()
    p.text = text
    p.font.size = Pt(font_size)
    p.font.color.rgb = color
    p.font.bold = bold
    p.level = level
    p.space_before = Pt(4)
    return p


# ═══════════════════════════════════════════════════════════════
#  SLIDE 1: Project Overview
# ═══════════════════════════════════════════════════════════════
slide1 = prs.slides.add_slide(prs.slide_layouts[6])  # blank
add_bg(slide1, BG_DARK)

# Title
add_text_box(slide1, 0.5, 0.3, 12, 1.0,
             "Surface Laptop Robot — Dual-Arm Autonomous Testing",
             font_size=32, color=ACCENT, bold=True, alignment=PP_ALIGN.LEFT)

# Subtitle
add_text_box(slide1, 0.5, 1.1, 8, 0.5,
             "AI-Driven Physical Device Interaction for Surface Quality Assurance",
             font_size=16, color=LIGHT_GRAY)

# Left column — project description
tf1 = add_text_box(slide1, 0.5, 1.8, 5.5, 5.0,
                   "Project Overview", font_size=22, color=ORANGE, bold=True)
add_bullet(tf1, "Dual myCobot 280 robot arms physically press keyboard keys", color=WHITE)
add_bullet(tf1, "and interact with touchpad on Surface laptop DUT", color=WHITE)
add_bullet(tf1, "")
add_bullet(tf1, "Key Components:", font_size=18, color=ORANGE, bold=True)
add_bullet(tf1, "• Ortler keyboard XML layout (78 keys with mm positions)", level=0)
add_bullet(tf1, "• Multi-camera system (3 USB + network cameras)")
add_bullet(tf1, "• Gambit REST API on DUT for HID stream detection")
add_bullet(tf1, "• Iterative position calibration with stream feedback")
add_bullet(tf1, "• Camera-to-robot hand-eye calibration (3-point anchor)")
add_bullet(tf1, "")
add_bullet(tf1, "Architecture:", font_size=18, color=ORANGE, bold=True)
add_bullet(tf1, "• Dev Machine → TCP → Raspberry Pi → Serial → myCobot 280")
add_bullet(tf1, "• Dev Machine → HTTP → Gambit (port 22133) → HID streams")
add_bullet(tf1, "• Keyboard stream: real-time key press detection")
add_bullet(tf1, "• Cursor stream: touchpad interaction verification")

# Right column — multi-camera image
try:
    # Try the composite camera image
    for img_path in [
        os.path.join(IMG_DIR, "camera_check", "all_now.jpg"),
        os.path.join(IMG_DIR, "camera_check", "composite_now.jpg"),
        os.path.join(IMG_DIR, "camera_check", "trio_now.jpg"),
    ]:
        if os.path.exists(img_path):
            slide1.shapes.add_picture(img_path, Inches(6.5), Inches(1.8), Inches(6.3), Inches(2.5))
            break
except:
    pass

# Stats box
tf_stats = add_text_box(slide1, 6.5, 4.5, 6.3, 2.5,
                        "System Stats", font_size=20, color=GREEN, bold=True)
add_bullet(tf_stats, "78 keys mapped from manufacturer XML layout", color=WHITE)
add_bullet(tf_stats, "43 single-character keys testable by robot", color=WHITE)
add_bullet(tf_stats, "32 keys calibrated & verified via HID stream", color=WHITE, bold=True)
add_bullet(tf_stats, "3 camera views: overhead, front, close-up", color=WHITE)
add_bullet(tf_stats, "Gambit plugins: Injection, Streams, Sensors", color=WHITE)
add_bullet(tf_stats, "Dual-arm: right (10.105.230.93) + left (10.105.230.94)", color=WHITE)

# ═══════════════════════════════════════════════════════════════
#  SLIDE 2: Achievements & Demo
# ═══════════════════════════════════════════════════════════════
slide2 = prs.slides.add_slide(prs.slide_layouts[6])  # blank
add_bg(slide2, BG_MED)

# Title
add_text_box(slide2, 0.5, 0.3, 12, 1.0,
             "Achievements & Calibration Pipeline",
             font_size=32, color=ACCENT, bold=True)

# Left column — achievements
tf2 = add_text_box(slide2, 0.5, 1.3, 5.5, 6.0,
                   "Key Achievements", font_size=22, color=GREEN, bold=True)
add_bullet(tf2, "✓ Camera annotation UI — 78 keys mapped in seconds", color=WHITE, bold=True)
add_bullet(tf2, "  Click 2 anchor keys → all keys auto-positioned via XML", color=LIGHT_GRAY)
add_bullet(tf2, "")
add_bullet(tf2, "✓ Gambit HID keyboard stream integration", color=WHITE, bold=True)
add_bullet(tf2, "  Real-time physical key press detection — no Notepad", color=LIGHT_GRAY)
add_bullet(tf2, "")
add_bullet(tf2, "✓ Iterative calibration with correction learning", color=WHITE, bold=True)
add_bullet(tf2, "  Press → detect → correct → verify 3x → save", color=LIGHT_GRAY)
add_bullet(tf2, "")
add_bullet(tf2, "✓ Dual-arm coordination", color=WHITE, bold=True)
add_bullet(tf2, "  Right arm: right-side keys (i,k,l,o,p,8,9...)", color=LIGHT_GRAY)
add_bullet(tf2, "  Left arm: left-side keys (a,s,d,f,g,q,w,e...)", color=LIGHT_GRAY)
add_bullet(tf2, "")
add_bullet(tf2, "✓ Touchpad localization from keyboard XML offset", color=WHITE, bold=True)
add_bullet(tf2, "  TP center = keyboard origin + (139.5mm, 155mm)", color=LIGHT_GRAY)
add_bullet(tf2, "")
add_bullet(tf2, "✓ Multi-camera GIF recording for demos", color=WHITE, bold=True)

# Right side — annotate screenshot
try:
    # Use the robot_dance GIF frame or annotation screenshot
    for img_path in [
        os.path.join(IMG_DIR, "camera_check", "all_now.jpg"),
    ]:
        if os.path.exists(img_path):
            slide2.shapes.add_picture(img_path, Inches(6.5), Inches(1.3), Inches(6.3), Inches(2.5))
            break
except:
    pass

# Pipeline diagram (text-based)
tf_pipe = add_text_box(slide2, 6.5, 4.0, 6.3, 3.2,
                       "Calibration Pipeline", font_size=20, color=ORANGE, bold=True)
add_bullet(tf_pipe, "1. annotate_keys.py — camera + XML → 78 key positions", color=WHITE)
add_bullet(tf_pipe, "2. Teach 3 anchor keys per arm (drag & record)", color=WHITE)
add_bullet(tf_pipe, "3. Fit affine: keyboard mm → robot coordinates", color=WHITE)
add_bullet(tf_pipe, "4. Predict all key positions from transform", color=WHITE)
add_bullet(tf_pipe, "5. Verify each via Gambit /streams/keyboard", color=WHITE)
add_bullet(tf_pipe, "6. Correct iteratively using XML key spacing", color=WHITE)
add_bullet(tf_pipe, "7. Save to learned_corrections.json", color=WHITE)
add_bullet(tf_pipe, "")
add_bullet(tf_pipe, "Goal: 100% key accuracy → type full paragraphs", color=GREEN, bold=True)

# Save
prs.save(OUTPUT)
print(f"Slide deck saved: {OUTPUT}")
